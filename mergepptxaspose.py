import urllib3
import changeWord
import asposeslidescloud
import shutil
from asposeslidescloud.apis.slides_api import SlidesApi
from asposeslidescloud.apis.slides_async_api import SlidesAsyncApi
from asposeslidescloud.models import *
import platform
import os
import collections
import collections.abc
from pptx import Presentation
import requests
import json
from urllib3.exceptions import InsecureRequestWarning
import time
from datetime import datetime, timedelta
# Suppress only the single warning from urllib3 needed.
requests.packages.urllib3.disable_warnings(category=InsecureRequestWarning)
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

def getfile_insensitive(paths):
    for path, subdirs, files in os.walk("/root/Dropbox/"):
        for name in files:
            #print(os.path.join(path, name).lower())
            if (os.path.join(path, name).lower() == paths.lower()):
                return (os.path.join(path, name))


def makeIntoList(y, date):
    answer = []
    try:
        response = requests.get('https://stmarkapi.com:5000/bishop?date=' + str(date), verify=False)
        bishop = json.loads(response.text)
        bishop = bishop["bishop"]
    except:
        bishop =""



    for i in y:
        if (i != "Ocassion" and i != "Season" and i != "Sunday" and i != "verb"):
            if (type(y[i]) is list):
                for l in y[i]:
                    l = l.replace('powerpoints', 'PowerPoints')
                    answer.append(l)
            else:
                if(y[i] == "alternate"):
                    y[i] = "PowerPoints/BackBone/AnotherLitanyOftheGospel.pptx"
                
                if(y[i] == "standard"):
                    y[i] = "PowerPoints/BackBone/litanyofthegospel.pptx"
                
                if(y[i] == "no"):
                    y[i] = ""

                if(i == "vespersOGodHaveMercy" and "Feast of the Cross" in y["vespersVerseofTheCymbals"]):
                    answer.append(y[i])
                    answer.append("PowerPoints/Feast of the Cross/Hymn of Constantine.pptx")
                    answer.append("PowerPoints/Feast of the Cross/Expositions.pptx")
                    y[i] = ""

                if(i == "vespers5ShortLitanies" and y[i] == "yes"):
                   y[i] = "PowerPoints/BackBone/5ShortLitanies.pptx" 
                   
                if(i == "matins5ShortLitanies" and y[i] == "yes"):
                    y[i] = "PowerPoints/BackBone/5ShortLitanies.pptx" 
                
                if(i == "Liturgy3GreatLitanies" and y[i] == "yes"):
                    y[i] = "PowerPoints/BackBone/threeGreatLitanies.pptx" 
                
                if(i == "rejoiceOMary" and y[i] == "yes"):
                    y[i] = "PowerPoints/Liturgy/Rejoice O Mary.pptx"
                    
                if(i == "OLordofHosts" and y[i] == "yes"):
                    y[i] = "PowerPoints/Liturgy/O Lord of Hosts.pptx"
                    y["theCherubim"] = ""

                if(i == "anaphora" and y[i] == "basil"):
                    y[i] = "PowerPoints/Liturgy/Anaphora - Basil.pptx"
    
                if(i == "anaphora" and y[i] == "gregory"):
                    y[i] = "PowerPoints/Liturgy/Anaphora - Gregorian.pptx"
  
                if(i == "agiosLiturgy" and y[i] == "basil"):
                    y[i] = "PowerPoints/Liturgy/Agios - Basil.pptx"
                
                if(i == "agiosLiturgy" and y[i] == "gregory"):
                    y[i] = "PowerPoints/Liturgy/Agios - Gregorian.pptx"  
                
                if(i == "instiution" and y[i] == "basil"):
                    y[i] = "PowerPoints/Liturgy/Institution - Basil.pptx"
                
                if(i == "instiution" and y[i] == "gregory"):
                    y[i] = "PowerPoints/Liturgy/Institution - Gregorian.pptx"  
                
                if(i == "yeahWeAskYou" and y[i] == "yes"):
                    y[i] = "PowerPoints/Liturgy/Yea we ask You.pptx"
                
                if(i == "jeNaiNan" and y[i] == "yes"):
                    y[i] = "PowerPoints/Liturgy/JeNaiNan.pptx"
                
                if(i == "healingToThesick" and y[i] == "yes"):
                    y[i] = "PowerPoints/Liturgy/Healing to the Sick.pptx"
                
                if(i == "Commemoration" and y[i] == "basil"):
                    y[i] = "PowerPoints/Liturgy/Commemoration - Basil.pptx"
                
                if(i == "Commemoration" and y[i] == "gregory"):
                    y[i] = "PowerPoints/Liturgy/Commemoration - Gregorian.pptx"  
                
                if(i == "postCommemoration" and y[i] == "basil"):
                    y[i] = "PowerPoints/Liturgy/Post Commemoration - Basil.pptx"
                
                if(i == "postCommemoration" and y[i] == "gregory"):
                    y[i] = "PowerPoints/Liturgy/Post Commemoration - Gregorian.pptx"  

                if (i == "Oblations" and is_within_four_sundays(str(date))):
                    answer.append("PowerPoints/Liturgy/Litany of the Leader.pptx")

                if(i == "prefaceToTheFraction" and y[i] == "basil"):
                    y[i] = "PowerPoints/Liturgy/Preface - Basil.pptx"
                
                if(i == "prefaceToTheFraction" and y[i] == "gregory"):
                    y[i] = "PowerPoints/Liturgy/Preface - Gregorian.pptx"  

                if(bishop == "yes"):
                    if(i == "vespersPrayerofThanksgiving"):
                        y[i] = "PowerPoints/BackBone/PrayerOfThanksgivingBishopVespers.pptx"

                    if(i == "vespers5ShortLitanies"):
                        y[i] = "PowerPoints/BackBone/5ShortLitanies.pptx"

                    if(i == "vespersConclusion"):
                        y[i] = "PowerPoints/BackBone/bishopConcludingHymn.pptx"

                    if(i == "OfferingThanksgiving"):
                        y[i] = "PowerPoints/BackBone/OfferingPrayerOfThanksgivingBishop.pptx"
                   
                    if(i == "hymnofIntercessions"):
                        answer.append(y[i])
                        answer.append("PowerPoints/BackBone/ResponseTothePaulineBishop.pptx")
                        y[i] = ""

                    if(i == "Liturgy3GreatLitanies"):
                        y[i] = "PowerPoints/BackBone/threeGreatLitanies.pptx" 
                    
                    if(i == "finalConclusion2"):
                        y[i] = "PowerPoints/BackBone/bishopConcludingHymn.pptx"
                    
                    if(i == "vespersDoxolgiesConcl"):
                        answer.append("PowerPoints/BackBone/BishopDoxology.pptx")
                        
                    if(i == "LiturgyPsalmResonse"):
                        answer.append("PowerPoints/BackBone/Psalm Trailer for Bishop.pptx")

                    if(i == "transitionSlide"):
                        answer.append("PowerPoints/BackBone/BishopEntryHymn.pptx")

                if (y[i] != ""):
                    answer.append(y[i])

    return (answer)


def merge(finishedList,date):
    import datetime
    today = datetime.date.today()
    currentDate = today.strftime("%Y-%m-%d")

    if ("Windows" in platform.platform()):
        path = "C:/Users/Mina Hanna/DropBox/"
    if (("Linux" in platform.platform())):
        path = "/root/Dropbox/"

    try:
        log = open(path + 'PowerPoints/configs/Logs/log ' + currentDate + '.txt', 'a')
    except:
        log = open(path + 'PowerPoints/configs/Logs/log ' + currentDate + '.txt', 'w')
    
    print(len(finishedList))


    presentaionsArray = []
    #finishedList.append("PowerPoints/communion.pptx")
    pptxLengths = {}
    totalBeforeCommunion = 0
    count = 0
    atMenu = False
    index = 0
    listTosend = []
    for i in range(0, len(finishedList), 10):
        
        try:
            log = open(path + 'PowerPoints/configs/Logs/log ' + currentDate + '.txt', 'a')
        except:
            log = open(path + 'PowerPoints/configs/Logs/log ' + currentDate + '.txt', 'w')
        
        files = []
        filesToremove = []
        x = 10
        if(len(finishedList[i+10:i+20]) < 2):
            x = 11
        for k in finishedList[i:i+x]:
            
            if ("today.pptx" in path + k):
                filesToremove.append(path+k)
            try:
                with open(path + k, "rb") as file_stream:
                    files.append(file_stream.read())

                    now = datetime.datetime.now()
                    timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
                    log.write(timestamp + " " + path + k + "\n")

                    print(path + k)
                    listTosend.append(path + k)

                    if(finishedList.index(k,index) > finishedList.index("PowerPoints/BackBone/communionMenuTemplate.pptx") and finishedList.index(k,index) < finishedList.index("PowerPoints/BackBone/finalConclusion1.pptx")):
                        pptxLengths[str(k.split("/")[-1].split(".")[0])] = (len(Presentation(path + k).slides))

                    if(finishedList.index(k,index) < finishedList.index("PowerPoints/BackBone/communionMenuTemplate.pptx") and atMenu == False):
                        totalBeforeCommunion = totalBeforeCommunion + (len(Presentation(path + k).slides))
                    if(k == "PowerPoints/BackBone/communionMenuTemplate.pptx"):
                        atMenu = True
                        #totalBeforeCommunion = totalBeforeCommunion - 1

                    
            except:
                try:
                    with open(getfile_insensitive(path + k), "rb") as file_stream:
                        files.append(file_stream.read())

                        now = datetime.datetime.now()
                        timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
                        log.write(timestamp + " " + path + k + "\n")

                        print(path + k)
                        listTosend.append(getfile_insensitive(path + k))

                        if(finishedList.index(k,index) > finishedList.index("PowerPoints/BackBone/communionMenuTemplate.pptx") and finishedList.index(k,index) < finishedList.index("PowerPoints/BackBone/finalConclusion1.pptx")):
                            pptxLengths[str(k.split("/")[-1].split(".")[0])] = (len(Presentation(getfile_insensitive(path + k)).slides))
                   
                        if(finishedList.index(k,index) < finishedList.index("PowerPoints/BackBone/communionMenuTemplate.pptx") and atMenu == False):
                            totalBeforeCommunion = totalBeforeCommunion + (len(Presentation(getfile_insensitive(path + k)).slides))
                        if(k == "PowerPoints/BackBone/communionMenuTemplate.pptx"):
                            atMenu = True
                except:
                    pass
            #print("Before Communion" ,  totalBeforeCommunion)
            index+=1
        
        print("uploading....")
        slides_api = SlidesApi(
            None, "2d3b1ec8-738b-4467-915f-af02913aa7fa", "1047551018f0feaacf4296fa054d7d97")
    
        # slides_api.merge_and_save_online(
        #     str(count) + ".pptx", files, None, "internal")

        presentation = PresentationToMerge()
        presentation.path = str(count) + ".pptx"
        presentation.source = "Storage"

        presentaionsArray.append(presentation)

        for i in filesToremove: 
            try:
                os.remove(i)
            except:
                pass

        count += 1
        if x == 11:
            break


    
    '''
    presentation = PresentationToMerge()
    presentation.path = "communion.pptx"
    presentation.source = "Storage"

    presentaionsArray.append(presentation)
    '''

    mergeSpringBoot(listTosend)
    # request = OrderedMergeRequest()
    # request.presentations = presentaionsArray
    # # response = slides_api.merge_and_save_online(
    # #     "MyPresentation.pptx", None,  request, "internal")
    
    # slides_asyncapi = SlidesAsyncApi(
    #         None, "2d3b1ec8-738b-4467-915f-af02913aa7fa", "1047551018f0feaacf4296fa054d7d97")

    # operation_id = slides_asyncapi.start_merge_and_save(out_path="MyPresentation.pptx", files=None,  request=request, storage="internal")

    # while True:
    #     time.sleep(2)
    #     operation = slides_asyncapi.get_operation_status(operation_id)
    #     print(f"Current operation status: { operation.status }")
    #     if operation.status == 'Started':
    #         if operation.progress != None:
    #             print(f"Operation is in progress. Merged { operation.progress.step_index } of { operation.progress.step_count }.")
    #             timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
    #             log.write(timestamp + " Operation is in progress. Merged " + str(operation.progress.step_index) + " of " + str(operation.progress.step_count) + ".\n")
    #     elif operation.status == 'Canceled':
    #         break
    #     elif operation.status == 'Failed':
    #         print(operation.error)
    #         timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
    #         log.write(timestamp + " " + str(operation.error) + "\n")
    #         break
    #     elif operation.status == 'Finished': 
    #         #result_path = slides_asyncapi.get_operation_result(operation_id)
    #         print("The merged document was Finished")
    #         timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
    #         log.write(timestamp + " The merged document was Finished\n")
    #         break
    

    # time.sleep(10)

    # try:
    #     slides_api.delete_unused_master_slides("MyPresentation.pptx", True)
    # except:
    #     pass
    
    result_path = path + "PowerPoints/result1.pptx"
    # temp_path = slides_api.download_file("MyPresentation.pptx")
    # shutil.copyfile(temp_path, result_path)

    print(pptxLengths)

    if(len(pptxLengths.keys()) > 0):
        from makeCommunionpptx import makePPT
        makePPT(result_path, pptxLengths, totalBeforeCommunion)

    response = requests.get('https://stmarkapi.com:8080/verb/?date='+date , verify=False)
            
    verb = response.text
    changeWord.insertChange(result_path,verb,"#SEASON#")
    changeWord.insertChange(result_path," ","#PLACEHOLDER#")

    postResponse = requests.get('https://stmarkapi.com:8080/pptname?date=' + date , verify=False)
    pptName = json.loads(postResponse.text)["pptName"]
    

    try:
        shutil.copyfile(result_path, "/root/" + pptName)
    except:
        os.makedirs("/root/" + "/".join(pptName.split("/")[:-1]))
        shutil.copyfile(result_path, "/root/" + pptName)

    path = pptName.replace("St. Mark PPT's" , "Current Liturgy PPT")
    print('complete')
    now = datetime.datetime.now()
    timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
    log.write(timestamp + " Complete\n")
    log.close()

def mergeCommunion(finishedList):
    
    platform.platform()
    if ("Windows" in platform.platform()):
        path = "C:/Users/minah/DropBox/"
    if (("Linux" in platform.platform())):
        path = "/root/Dropbox/"

    presentaionsArray = []

    count = 0
    for i in range(0, len(finishedList), 10):
        files = []
        filesToremove = []
        pptxLengths = {}
        for k in finishedList[i:i+10]:
            print(path + k)
            if ("today.pptx" in path + k):
                filesToremove.append(path+k)
            try:
                with open(path + k, "rb") as file_stream:
                    files.append(file_stream.read())
                    #print(str(k.split("/")[-1].split(".")[0]))
                    if(finishedList.index(k) > finishedList.index("PowerPoints/BackBone/communionMenuTemplate.pptx") and finishedList.index(k) < finishedList.index("PowerPoints/BackBone/finalConclusion1.pptx")):
                        pptxLengths[str(k.split("/")[-1].split(".")[0])] = (len(Presentation(path + k).slides))
            except:
                #print("here", getfile_insensitive(path+i))
                try:
                    with open(getfile_insensitive(path + k), "rb") as file_stream:
                        files.append(file_stream.read())
                        #print(str(k.split("/")[-1].split(".")[0]))
                        if(k!= finishedList[0] and k != finishedList[-1] and k != finishedList[-2] and k != finishedList[-3] and k != finishedList[-4]):
                            pptxLengths[str(k.split("/")[-1].split(".")[0])] = (len(Presentation(getfile_insensitive(path + k)).slides))
                except:
                    pass

        print("uploading....")
        print(pptxLengths)
        slides_api = SlidesApi(
            None, "2d3b1ec8-738b-4467-915f-af02913aa7fa", "1047551018f0feaacf4296fa054d7d97")
        slides_api.merge_and_save_online("communion.pptx", files, None, "internal")

    result_path = path + "PowerPoints/communion.pptx"
    temp_path = slides_api.download_file("communion.pptx", "internal")
    shutil.copyfile(temp_path, result_path)

    print('creating menu slide')
    
    if(len(pptxLengths.keys()) > 0):
        from makeCommunionpptx import makePPT
        makePPT(result_path, pptxLengths)

    print('complete')


def mergeSpringBoot(file_paths):
    
# Define the URL of the endpoint
    url = "https://stmarkapi.com:8080/mergePowerpoints"

    # Define the list of file paths

    # Convert the list to a JSON payload
    payload = json.dumps(file_paths)

    # Set headers to indicate JSON content
    headers = {
        "Content-Type": "application/json"
    }

    # Send the POST request
    response = requests.post(url, data=payload, headers=headers, verify=False)

    # Print the response from the server
    print("Status Code:", response.status_code)
    print("Response Body:", response.text)


def is_election_year(year):
    """Check if the given year is a US general election year."""
    return (year - 2020) % 4 == 0

def get_election_day(year):
    """Get the date of the general election day in a given year."""
    if not is_election_year(year):
        return None
    
    # Find the first Monday of November.
    november_first = datetime(year, 11, 1)
    first_monday = november_first + timedelta(days=(7 - november_first.weekday()) % 7)
    # Election Day is the day after the first Monday (i.e., Tuesday).
    election_day = first_monday + timedelta(days=1)
    return election_day

def is_within_four_sundays(date_str):
    """Check if the given date is within the 4 Sundays leading up to Election Day."""
    try:
        date = datetime.strptime(date_str, "%Y-%m-%d")
    except ValueError:
        raise ValueError("The date format should be YYYY-MM-DD.")
    
    year = date.year
    election_day = get_election_day(year)
    
    if election_day is None:
        return False  # Not an election year
    
    # Find the 4 Sundays before the election
    sunday_before_election = election_day - timedelta(days=(election_day.weekday() + 1))
    four_sundays_before_election = sunday_before_election - timedelta(weeks=3)
    
    # Check if the given date is within this range
    return four_sundays_before_election <= date <= sunday_before_election