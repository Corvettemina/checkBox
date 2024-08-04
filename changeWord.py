import collections
import collections.abc
from pptx import Presentation
import platform
from pptx.exc import InvalidXmlError

def insertChange(input_pptx, replaceString):
    # Determine the base path based on the operating system
    if "Windows" in platform.system():
        base_path = "C:/Users/Mina Hanna/DropBox/"
    elif "Linux" in platform.system():
        base_path = "/root/Dropbox/"
    else:
        raise OSError("Unsupported operating system")

    prs = Presentation(input_pptx)
    testString = "#SEASON#"

    # Get all shapes in all slides
    shapes = [shape for slide in prs.slides for shape in slide.shapes]

    replaces = {testString: replaceString}

    for shape in shapes:
        if not shape.has_text_frame:
            continue
        
        text_frame = shape.text_frame
        if text_frame is None:
            continue
        
        for match, replacement in replaces.items():
            try:
                if any(match in paragraph.text for paragraph in text_frame.paragraphs):
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(match, replacement)
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                paragraph._p.remove(run._r)
                        if paragraph.runs:
                            paragraph.runs[0].text = whole_text
            except InvalidXmlError:
                print(f"InvalidXmlError encountered in shape. Skipping this shape.")
                continue

    prs.save(input_pptx)
    return input_pptx

def main():
    input_pptx = "C:/Users/Mina Hanna/DropBox/PowerPoints/result1.pptx"
    replaceString = "have come"
    insertChange(input_pptx, replaceString)

if __name__ == "__main__":
    main()
