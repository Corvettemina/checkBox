import collections 
import collections.abc
from pptx import Presentation
import os

prs1 = Presentation("KihakVirginDoxology.pptx")
prs2 = Presentation("KhiakGabrilDoxology.pptx")

for slide in prs2.slides:
    try:
        sl = prs1.slides.add_slide(prs1.slide_layouts[1])
        sl.shapes.title.text = slide.shapes.title.text
        sl.placeholders[1].text = slide.placeholders[1].text
    except:
        pass

prs1.save("test.pptx")