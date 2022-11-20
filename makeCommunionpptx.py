import collections
import collections.abc

from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def makePPT(result_path, indexes, beforeCommunion):
    #readings = getReadingsDictionary(year, month, day)
    print(indexes)
    prs = Presentation(result_path)
    beforeCommunion
    slide = prs.slides[beforeCommunion]
    currentIndex = beforeCommunion + 1 
    verticalPostion = .48
    horizonatlPosition = .48 
    for i in indexes.keys():
        if(verticalPostion > 5.5):
            verticalPostion = 1.48
            horizonatlPosition = horizonatlPosition + 3.62

        text = i    
        blessed = slide.shapes.add_textbox(
            Inches(horizonatlPosition), Inches(verticalPostion), Inches(2.57), Inches(.54))
        tf = blessed.text_frame

        p = tf.paragraphs[0]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = True
        blessed.fill.solid()
        blessed.fill.fore_color.rgb = RGBColor(255, 255, 255)
        click_action = blessed.click_action
        click_action.target_slide = prs.slides[currentIndex]
        click_action.action
        line = blessed.line
        line.color.rgb = RGBColor(79, 129, 189)
        line.width = Pt(2)

        run = p.add_run()

        run.text = text
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(24)

        verticalPostion = verticalPostion + 1
        currentIndex = currentIndex + indexes[i]
        print("current index" ,currentIndex)
        #print(verticalPostion)

    for i in range(beforeCommunion, int(currentIndex)):
        slide = prs.slides[i]
        text = "Menu"
        menu = slide.shapes.add_textbox(
            Inches(.38), Inches(6.62), Inches(1.08), Inches(.54))
        tf = menu.text_frame

        p = tf.paragraphs[0]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = True
        menu.fill.solid()
        menu.fill.fore_color.rgb = RGBColor(255, 255, 255)
        click_action = menu.click_action
        click_action.target_slide = prs.slides[beforeCommunion]
        click_action.action
        line = menu.line
        line.color.rgb = RGBColor(79, 129, 189)
        line.width = Pt(2)

        run = p.add_run()

        run.text = text
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(24)

        text = "Conclusion"
        conclusion = slide.shapes.add_textbox(
            Inches(11.2), Inches(6.62), Inches(1.76), Inches(.54))
        tf = conclusion.text_frame

        p = tf.paragraphs[0]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = True
        conclusion.fill.solid()
        conclusion.fill.fore_color.rgb = RGBColor(255, 255, 255)
        click_action = conclusion.click_action
        click_action.target_slide = prs.slides[currentIndex]
        click_action.action
        line = conclusion.line
        line.color.rgb = RGBColor(79, 129, 189)
        line.width = Pt(2)

        run = p.add_run()

        run.text = text
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(24)

    prs.save(result_path)

    #import os
    #s.startfile("communion2.pptx")
#dictionary = {'standardPsalm150': 6, 'Fraction to the Father for Advent and the Nativity': 3}
#path = "C:/Users/Mina Hanna/Dropbox/PowerPoints/result1.pptx"
#makePPT(path, dictionary, 457)